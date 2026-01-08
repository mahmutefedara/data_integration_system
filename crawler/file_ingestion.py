import os
import tempfile
from urllib.parse import urlparse

import fitz
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from utils import hash_url


def _ext(url: str) -> str:
    return os.path.splitext(urlparse(url).path)[1].lower()


def extract_text_from_file(path: str, ext: str) -> str:
    try:
        if ext == ".pdf":
            doc = fitz.open(path)
            try:
                return "\n".join(page.get_text() for page in doc)
            finally:
                doc.close()

        if ext == ".docx": # .doc'u ayırdık çünkü python-docx .doc okuyamaz
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs)

        if ext == ".doc":
            # NOT: python-docx .doc (eski format) desteklemez.
            # Şimdilik hata vermemesi için boş dönüyoruz.
            # Gerçek çözüm için 'tika' veya 'antiword' kütüphanesi gerekir.
            return "[Legacy .doc file - extraction not supported with current library]"

        if ext in (".ppt", ".pptx"):
            prs = Presentation(path)
            out = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        t = (shape.text or "").strip()
                        if t:
                            out.append(t)
            return "\n".join(out)

        if ext in (".xls", ".xlsx"):
            wb = load_workbook(path, data_only=True)
            out = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    vals = [str(c) for c in row if c is not None and str(c).strip()]
                    if vals:
                        out.append(" ".join(vals))
            return "\n".join(out)

        if ext == ".txt":
            with open(path, "r", encoding="utf-8", errors="ignore") as f:
                return f.read()

    except Exception as e:
        # Herhangi bir dosya okuma hatasında crawler çökmesin diye hata mesajı dönüyoruz
        return f"[Error extracting {ext} file: {str(e)}]"

    return ""


async def download_extract_delete(fetcher, url: str, *, max_bytes: int | None = None):
    ext = _ext(url) or ".bin"
    tmp_dir = tempfile.gettempdir()
    tmp_path = os.path.join(tmp_dir, f"crawl_{hash_url(url)}{ext}")

    try:
        # Crawler'dan fetcher ve url'nin doğru sırayla geldiğinden emin ol
        data, ctype = await fetcher.fetch(url)
        if not data:
            return "", {"ext": ext, "size": 0}, ctype

        if max_bytes is not None and len(data) > max_bytes:
            return "", {"ext": ext, "size": len(data), "skipped_too_large": True}, ctype

        with open(tmp_path, "wb") as f:
            f.write(data)

        # Hata kontrolünü burada yapıyoruz
        text = extract_text_from_file(tmp_path, ext)
        meta = {"ext": ext, "size": len(data)}
        return text, meta, ctype

    except Exception as e:
        # İndirme veya yazma aşamasındaki hatalar için
        return "", {"ext": ext, "size": 0, "error": str(e)}, None

    finally:
        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except Exception:
            pass