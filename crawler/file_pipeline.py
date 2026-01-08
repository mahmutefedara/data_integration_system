import os
import tempfile
from urllib.parse import urlparse

import fitz  # pymupdf
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook

from utils import hash_url


def _ext(url: str) -> str:
    return os.path.splitext(urlparse(url).path)[1].lower()


def extract_text_from_file(path: str, ext: str) -> str:
    if ext == ".pdf":
        doc = fitz.open(path)
        try:
            return "\n".join(page.get_text() for page in doc)
        finally:
            doc.close()

    if ext in (".doc", ".docx"):
        doc = Document(path)
        return "\n".join(p.text for p in doc.paragraphs)

    if ext in (".ppt", ".pptx"):
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    t = (shape.text or "").strip()
                    if t:
                        out.append(t)
        return "\n".join(out)

    if ext in (".xls", ".xlsx"):
        wb = load_workbook(path, data_only=True)
        out = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                vals = [str(c) for c in row if c is not None and str(c).strip()]
                if vals:
                    out.append(" ".join(vals))
        return "\n".join(out)

    if ext == ".txt":
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()

    return ""


async def download_extract_delete(fetcher, url: str, *, max_bytes: int = 40_000_000):
    """
    Returns: (text, meta, content_type)
      meta = {"ext": ext, "source_bytes": int}
    """
    ext = _ext(url) or ".bin"
    tmp_dir = tempfile.gettempdir()
    tmp_path = os.path.join(tmp_dir, f"crawl_{hash_url(url)}{ext}")

    data = b""
    ctype = ""

    try:
        data, ctype = await fetcher.fetch(url)
        if not data:
            return "", {"ext": ext, "source_bytes": 0}, ctype

        if len(data) > max_bytes:
            return "", {"ext": ext, "source_bytes": len(data), "skipped": "too_large"}, ctype

        with open(tmp_path, "wb") as f:
            f.write(data)

        text = extract_text_from_file(tmp_path, ext)
        meta = {"ext": ext, "source_bytes": len(data)}
        return text, meta, ctype

    finally:
        try:
            if os.path.exists(tmp_path):
                os.remove(tmp_path)
        except Exception:
            pass
